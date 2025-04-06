# create_report.py
from pptx import Presentation
from pptx.util import Inches

def create_report(insights, chart_path="performance_chart.png"):
    """
    Crée un rapport PowerPoint en utilisant les données extraites et les graphiques générés.
    
    Args:
        insights (dict): Les informations extraites concernant le fonds.
        chart_path (str): Le chemin de l'image à inclure dans le rapport.
    """
    prs = Presentation()

    # Titre de la présentation
    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_1.shapes.title
    subtitle = slide_1.placeholders[1]
    title.text = "Rapport de Performance du Fonds"
    subtitle.text = "Généré automatiquement à partir du système Q&A"

    # Aperçu des performances
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_2.shapes.title.text = "Aperçu des Performances"
    slide_2.shapes.placeholders[1].text = f"Performance du fonds : {insights['performance']}"

    # Ajouter un graphique de performance
    slide_3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide_3.shapes.title.text = "Graphique de Performance"
    slide_3.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(8))

    # Conformité et autres métriques
    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_4.shapes.title.text = "Métriques de Conformité"
    slide_4.shapes.placeholders[1].text = f"Conformité : {insights['compliance']}"

    # Sauvegarder la présentation
    prs.save("rapport_fonds.pptx")
